#!/usr/bin/env python3
"""
Generate PowerPoint presentation from PROJECT_PRESENTATION.md
Creates slides with placeholders for screenshots
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def parse_markdown_slides(md_file):
    """Parse markdown file into slide content"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = []
    slide_pattern = r'## Slide (\d+): (.+?)\n\n(.*?)(?=\n---\n|## Slide|\Z)'
    
    matches = re.finditer(slide_pattern, content, re.DOTALL)
    for match in matches:
        slide_num = int(match.group(1))
        title = match.group(2)
        body = match.group(3).strip()
        
        # Parse bullet points
        bullets = []
        current_section = None
        
        for line in body.split('\n'):
            line = line.strip()
            if not line:
                continue
            
            # Check for section headers (###)
            if line.startswith('### **'):
                current_section = line.replace('### **', '').replace('**', '')
                bullets.append(('section', current_section))
            elif line.startswith('- **') or line.startswith('* **'):
                # Bullet point with bold
                bullet_text = line[2:].strip()
                bullets.append(('bullet', bullet_text))
            elif line.startswith('- ') or line.startswith('* '):
                # Regular bullet point
                bullets.append(('bullet', line[2:].strip()))
            elif line.startswith('```'):
                # Code block - skip for now or add as special
                continue
            elif line and not line.startswith('#'):
                # Regular text
                bullets.append(('text', line))
        
        slides.append({
            'number': slide_num,
            'title': title,
            'content': bullets
        })
    
    return slides

def add_text_to_shape(shape, text, is_bold=False, font_size=14):
    """Add formatted text to a shape"""
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    
    # Split text by ** markers for bold formatting
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            # Bold text
            run = p.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(0, 0, 0)
        elif part.strip():
            # Regular text
            run = p.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(0, 0, 0)
            if is_bold:
                run.font.bold = True

def create_slide(prs, slide_data):
    """Create a slide from slide data"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.9))
    title_frame = title_shape.text_frame
    title_frame.text = slide_data['title']
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add title underline using a rectangle
    underline = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.05))
    underline.fill.solid()
    underline.fill.fore_color.rgb = RGBColor(0, 51, 102)
    underline.line.fill.background()
    
    # Content area
    content_top = Inches(1.3)
    content_left = Inches(0.7)
    content_width = Inches(8.6)
    content_height = Inches(5.0)
    
    # Check if we need screenshot placeholder or diagram
    screenshot_keywords = ['visualization', 'CFG', 'graph', 'interconnected', 'call graph', 'taint analysis', 'visualizer']
    needs_screenshot = any(kw in slide_data['title'].lower() for kw in screenshot_keywords)
    
    # Add screenshot placeholders for key visualization slides
    screenshot_slides = [3, 32, 33, 37]  # CFG Generation, CFGVisualizer slides, Visualization Features
    
    # Add diagram placeholders for architecture/flow slides
    diagram_slides = [2, 7, 34, 35]  # Architecture, DataflowAnalyzer pipeline, IPA framework
    
    if slide_data['number'] in diagram_slides:
        # Add diagram placeholder
        diagram_shape = slide.shapes.add_textbox(content_left + Inches(0.5), content_top, content_width - Inches(1), Inches(3.5))
        diagram_frame = diagram_shape.text_frame
        diagram_frame.word_wrap = True
        
        # Determine diagram instructions based on slide
        if slide_data['number'] == 2:
            instructions = "[ARCHITECTURE DIAGRAM PLACEHOLDER]\n\nAdd diagram showing:\n- Analysis Pipeline (4 stages)\n- Key Components (5 components)\n- Data flow between components\n- CFG → Analysis → Visualization flow"
        elif slide_data['number'] == 7:
            instructions = "[ANALYSIS PIPELINE DIAGRAM PLACEHOLDER]\n\nAdd diagram showing:\n- 6-step pipeline\n- File parsing → Intra-procedural → Call graph\n- Inter-procedural → Visualization → State\n- Component interactions"
        elif slide_data['number'] == 34:
            instructions = "[IPA FRAMEWORK DIAGRAM PLACEHOLDER]\n\nAdd diagram showing:\n- Phase 1: Call Graph Construction\n- Phase 2: Advanced Call Graph Analysis\n- Component relationships\n- Data flow"
        elif slide_data['number'] == 35:
            instructions = "[IPA FRAMEWORK DIAGRAM PLACEHOLDER]\n\nAdd diagram showing:\n- Phase 3: Inter-Procedural Reaching Definitions\n- Phase 4: Parameter & Return Value Analysis\n- Component relationships\n- Data flow"
        else:
            instructions = "[DIAGRAM PLACEHOLDER]\n\nAdd diagram showing:\n- Component relationships\n- Data flow\n- Analysis pipeline"
        
        diagram_frame.text = instructions
        diagram_para = diagram_frame.paragraphs[0]
        diagram_para.font.size = Pt(12)
        diagram_para.font.italic = True
        diagram_para.font.color.rgb = RGBColor(150, 150, 150)
        diagram_para.alignment = PP_ALIGN.CENTER
        
        diagram_shape.line.color.rgb = RGBColor(200, 200, 200)
        diagram_shape.line.width = Pt(2)
        diagram_shape.fill.solid()
        diagram_shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
        
        # Text content below diagram
        text_shape = slide.shapes.add_textbox(content_left, content_top + Inches(3.8), content_width, Inches(2.8))
        text_frame = text_shape.text_frame
        text_frame.word_wrap = True
    elif needs_screenshot and slide_data['number'] in screenshot_slides:
        # Split layout: text on left, screenshot placeholder on right
        text_width = Inches(4.0)
        screenshot_left = Inches(5.3)
        screenshot_width = Inches(4.0)
        screenshot_height = Inches(3.5)
        
        # Text content
        text_shape = slide.shapes.add_textbox(content_left, content_top, text_width, content_height)
        text_frame = text_shape.text_frame
        text_frame.word_wrap = True
        
        # Screenshot placeholder
        screenshot_shape = slide.shapes.add_textbox(screenshot_left, content_top + Inches(0.2), screenshot_width, screenshot_height)
        screenshot_frame = screenshot_shape.text_frame
        screenshot_frame.word_wrap = True
        
        # Determine screenshot instructions based on slide
        if slide_data['number'] == 3:
            instructions = "CFG Generation Pipeline\n\nTake screenshot of:\n- cfg-exporter tool\n- CMake build output\n- JSON CFG output"
        elif slide_data['number'] == 32:
            instructions = "CFGVisualizer Overview\n\nTake screenshot of:\n- VS Code extension panel\n- Visualization tabs\n- CFG graph display"
        elif slide_data['number'] == 33:
            instructions = "CFGVisualizer Implementation\n\nTake screenshot of:\n- Interconnected CFG view\n- Edge type toggles\n- Multiple visualization tabs"
        elif slide_data['number'] == 37:
            instructions = "Visualization Features\n\nTake screenshot of:\n- Interactive CFG\n- Color-coded nodes\n- Edge type visualization"
        else:
            instructions = f"[SCREENSHOT PLACEHOLDER]\n\n{slide_data['title']}\n\nTake screenshot of:\n- Visualization panel\n- Relevant tabs\n- Analysis results"
        
        screenshot_frame.text = instructions
        screenshot_para = screenshot_frame.paragraphs[0]
        screenshot_para.font.size = Pt(12)
        screenshot_para.font.italic = True
        screenshot_para.font.color.rgb = RGBColor(128, 128, 128)
        screenshot_para.alignment = PP_ALIGN.CENTER
        
        # Add border to placeholder
        screenshot_shape.line.color.rgb = RGBColor(200, 200, 200)
        screenshot_shape.line.width = Pt(2)
        screenshot_shape.fill.solid()
        screenshot_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    else:
        # Full width text
        text_shape = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        text_frame = text_shape.text_frame
        text_frame.word_wrap = True
    
    # Add content
    for item_type, content in slide_data['content']:
        if item_type == 'section':
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = content
            run.font.bold = True
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0, 51, 102)
            p.space_after = Pt(6)
        elif item_type == 'bullet':
            p = text_frame.add_paragraph()
            p.level = 0
            p.space_before = Pt(3)
            # Parse bold markers in bullet text
            parts = re.split(r'(\*\*.*?\*\*)', content)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run()
                    run.text = part[2:-2]
                    run.font.bold = True
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                elif part.strip():
                    run = p.add_run()
                    run.text = part
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(0, 0, 0)
        elif item_type == 'text':
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = content
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(100, 100, 100)
            p.space_after = Pt(3)
    
    # Slide number
    slide_num_shape = slide.shapes.add_textbox(Inches(9.2), Inches(7.0), Inches(0.5), Inches(0.3))
    slide_num_frame = slide_num_shape.text_frame
    slide_num_frame.text = str(slide_data['number'])
    slide_num_para = slide_num_frame.paragraphs[0]
    slide_num_para.font.size = Pt(10)
    slide_num_para.font.color.rgb = RGBColor(150, 150, 150)
    slide_num_para.alignment = PP_ALIGN.RIGHT

def main():
    """Main function to generate presentation"""
    print("Generating PowerPoint presentation...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Parse markdown
    slides_data = parse_markdown_slides('PROJECT_PRESENTATION.md')
    
    print(f"Found {len(slides_data)} slides")
    
    # Create slides
    for slide_data in slides_data:
        print(f"Creating slide {slide_data['number']}: {slide_data['title']}")
        create_slide(prs, slide_data)
    
    # Save presentation
    output_file = 'C++_Dataflow_Analyzer_Presentation.pptx'
    prs.save(output_file)
    print(f"\nPresentation saved as: {output_file}")
    print(f"Total slides: {len(slides_data)}")

if __name__ == '__main__':
    main()

